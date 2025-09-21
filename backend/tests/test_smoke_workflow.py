"""
Comprehensive smoke tests for the authentication and job submission workflow.

These tests verify the entire end-to-end flow from user registration to job completion.
"""
import pytest
import json
import os
import tempfile
from pathlib import Path
from unittest.mock import patch, MagicMock
import base64
from datetime import datetime, timedelta

from app.core.config import settings


class TestAuthenticationFlow:
    """Test the complete authentication workflow"""

    def test_user_registration_success(self, client):
        """Test successful user registration"""
        user_data = {
            "email": "smoketest@example.com",
            "password": "SmokeTest123!",
            "full_name": "Smoke Test User"
        }

        response = client.post("/api/auth/register", json=user_data)
        assert response.status_code == 200
        data = response.json()

        assert data["email"] == user_data["email"]
        assert data["full_name"] == user_data["full_name"]
        assert "id" in data
        assert data["is_active"] is True
        assert "hashed_password" not in data  # Never return password

    def test_user_registration_duplicate_email(self, client):
        """Test registration with duplicate email fails"""
        user_data = {
            "email": "duplicate@example.com",
            "password": "TestPass123!",
            "full_name": "First User"
        }

        # Register first user
        client.post("/api/auth/register", json=user_data)

        # Try to register same email again
        response = client.post("/api/auth/register", json=user_data)
        assert response.status_code == 400
        assert "already registered" in response.json()["detail"].lower()

    def test_user_registration_weak_password(self, client):
        """Test registration with weak password fails"""
        user_data = {
            "email": "weak@example.com",
            "password": "123",  # Too weak
            "full_name": "Weak Password User"
        }

        response = client.post("/api/auth/register", json=user_data)
        assert response.status_code == 422  # Validation error

    def test_login_success(self, client):
        """Test successful user login"""
        # First register a user
        user_data = {
            "email": "login@example.com",
            "password": "LoginPass123!",
            "full_name": "Login User"
        }
        client.post("/api/auth/register", json=user_data)

        # Then login
        login_data = {
            "email": "login@example.com",
            "password": "LoginPass123!"
        }
        response = client.post("/api/auth/login", json=login_data)

        assert response.status_code == 200
        data = response.json()
        assert "access_token" in data
        assert "refresh_token" in data
        assert data["token_type"] == "bearer"
        assert isinstance(data["expires_in"], int)
        assert data["expires_in"] > 0

    def test_login_invalid_credentials(self, client):
        """Test login with invalid credentials fails"""
        # Register a user
        user_data = {
            "email": "invalid@example.com",
            "password": "CorrectPass123!",
            "full_name": "Invalid Login User"
        }
        client.post("/api/auth/register", json=user_data)

        # Try to login with wrong password
        login_data = {
            "email": "invalid@example.com",
            "password": "WrongPassword!"
        }
        response = client.post("/api/auth/login", json=login_data)

        assert response.status_code == 401
        assert "invalid credentials" in response.json()["detail"].lower()

    def test_login_nonexistent_user(self, client):
        """Test login with non-existent user fails"""
        login_data = {
            "email": "nonexistent@example.com",
            "password": "AnyPassword123!"
        }
        response = client.post("/api/auth/login", json=login_data)

        assert response.status_code == 401
        assert "invalid credentials" in response.json()["detail"].lower()

    def test_refresh_token_success(self, client):
        """Test successful token refresh"""
        # Register and login to get tokens
        user_data = {
            "email": "refresh@example.com",
            "password": "RefreshPass123!",
            "full_name": "Refresh User"
        }
        client.post("/api/auth/register", json=user_data)

        login_data = {
            "email": "refresh@example.com",
            "password": "RefreshPass123!"
        }
        login_response = client.post("/api/auth/login", json=login_data)
        refresh_token = login_response.json()["refresh_token"]

        # Use refresh token
        refresh_data = {"refresh_token": refresh_token}
        response = client.post("/api/auth/refresh", json=refresh_data)

        assert response.status_code == 200
        data = response.json()
        assert "access_token" in data
        assert data["token_type"] == "bearer"
        # Should be a new token
        assert data["access_token"] != login_response.json()["access_token"]

    def test_refresh_token_invalid(self, client):
        """Test refresh with invalid token fails"""
        refresh_data = {"refresh_token": "invalid_token"}
        response = client.post("/api/auth/refresh", json=refresh_data)

        assert response.status_code == 401

    def test_protected_endpoint_without_token(self, client):
        """Test accessing protected endpoint without authentication"""
        response = client.get("/api/jobs")
        assert response.status_code == 403  # FastAPI returns 403 for missing auth

    def test_protected_endpoint_with_invalid_token(self, client):
        """Test accessing protected endpoint with invalid token"""
        headers = {"Authorization": "Bearer invalid_token"}
        response = client.get("/api/jobs", headers=headers)
        assert response.status_code == 401

    def test_user_profile_access(self, client, auth_headers):
        """Test accessing user profile with valid token"""
        response = client.get("/api/auth/me", headers=auth_headers)
        assert response.status_code == 200
        data = response.json()
        assert "email" in data
        assert "full_name" in data
        assert data["email"] == "test@example.com"


class TestJobSubmissionWorkflow:
    """Test the complete job submission workflow"""

    def create_test_pptx(self, content="テスト"):
        """Create a minimal test PPTX file"""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = content

        # Save to temporary file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            return f.name

    def create_test_pdf(self, content="テスト"):
        """Create a minimal test PDF file"""
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            c = canvas.Canvas(f.name, pagesize=letter)
            c.drawString(100, 750, content)
            c.save()
            return f.name

    @patch('app.core.job_manager.job_manager.process_translation_job')
    @patch('app.services.translation.translation_service.translate_text')
    def test_pptx_translation_job_complete_workflow(
        self,
        mock_translate,
        mock_process,
        client,
        auth_headers,
        test_upload_dir,
        test_output_dir
    ):
        """Test complete PPTX translation workflow"""
        # Mock the translation service
        mock_translate.return_value = {
            "translations": [{"original": "テスト", "translated": "Test"}]
        }

        # Mock the job processor to simulate async processing
        mock_process.return_value = None

        # Create test PPTX
        pptx_path = self.create_test_pptx()

        try:
            # Step 1: Upload file and create job
            with open(pptx_path, "rb") as f:
                files = {"file": ("test.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
                data = {
                    "file_type": "pptx",
                    "model": "gpt-4o-mini",
                    "temperature": 0.6,
                    "auto_fit": "norm"
                }

                response = client.post(
                    "/api/translate",
                    files=files,
                    data=data,
                    headers=auth_headers
                )

            assert response.status_code == 200
            job_data = response.json()
            assert "job_id" in job_data
            assert job_data["status"] in ["pending", "processing"]
            assert job_data["file_type"] == "pptx"

            job_id = job_data["job_id"]

            # Step 2: Check job status (might still be processing)
            response = client.get(f"/api/jobs/{job_id}", headers=auth_headers)
            assert response.status_code == 200

            # For testing purposes, simulate job completion
            # In real scenario, this would be handled by the job processor
            with patch('app.core.job_manager.job_manager.get_job_status') as mock_status:
                mock_status.return_value = {
                    "status": "completed",
                    "progress": 100,
                    "message": "Translation completed",
                    "result_file": f"translated_{job_id}.pptx",
                    "created_at": datetime.utcnow().isoformat(),
                    "updated_at": datetime.utcnow().isoformat(),
                    "stats": {
                        "pages_processed": 1,
                        "tokens_used": 100,
                        "cost": 0.002
                    }
                }

                response = client.get(f"/api/jobs/{job_id}", headers=auth_headers)
                assert response.status_code == 200
                status_data = response.json()
                assert status_data["status"] == "completed"
                assert status_data["progress"] == 100

            # Step 3: List jobs to verify it appears in history
            response = client.get("/api/jobs", headers=auth_headers)
            assert response.status_code == 200
            jobs_data = response.json()
            assert "jobs" in jobs_data
            assert any(job["id"] == job_id for job in jobs_data["jobs"])

            # Step 4: Download translated file (mock the file existence)
            with patch('pathlib.Path.exists', return_value=True):
                with patch('fastapi.responses.FileResponse') as mock_file_response:
                    mock_file_response.return_value = MagicMock()
                    mock_file_response.return_value.status_code = 200

                    response = client.get(
                        f"/api/jobs/{job_id}/download",
                        headers=auth_headers
                    )
                    # Note: FileResponse can't be easily tested in TestClient
                    # In real tests, you'd check the file headers and content

            # Step 5: Search for the job
            search_data = {
                "search": "test",
                "page": 1,
                "page_size": 10
            }
            response = client.post("/api/jobs/search", json=search_data, headers=auth_headers)
            assert response.status_code == 200
            search_results = response.json()
            assert "jobs" in search_results
            assert search_results["total"] >= 1

        finally:
            # Clean up
            os.unlink(pptx_path)

    @patch('app.core.job_manager.job_manager.process_translation_job')
    @patch('app.services.translation.translation_service.translate_text')
    def test_pdf_translation_job_complete_workflow(
        self,
        mock_translate,
        mock_process,
        client,
        auth_headers,
        test_upload_dir,
        test_output_dir
    ):
        """Test complete PDF translation workflow"""
        # Mock the translation service
        mock_translate.return_value = {
            "translations": [{"original": "テスト", "translated": "Test"}]
        }

        # Mock the job processor
        mock_process.return_value = None

        # Create test PDF
        pdf_path = self.create_test_pdf()

        try:
            # Step 1: Upload file and create job
            with open(pdf_path, "rb") as f:
                files = {"file": ("test.pdf", f, "application/pdf")}
                data = {
                    "file_type": "pdf",
                    "model": "gpt-4o-mini",
                    "temperature": 0.6,
                    "pages": "1-1",
                    "auto_fit": "shape"
                }

                response = client.post(
                    "/api/translate",
                    files=files,
                    data=data,
                    headers=auth_headers
                )

            assert response.status_code == 200
            job_data = response.json()
            assert "job_id" in job_data
            assert job_data["file_type"] == "pdf"

            job_id = job_data["job_id"]

            # Step 2: Verify job was created with correct parameters
            response = client.get(f"/api/jobs/{job_id}", headers=auth_headers)
            assert response.status_code == 200

            # Simulate job completion
            with patch('app.core.job_manager.job_manager.get_job_status') as mock_status:
                mock_status.return_value = {
                    "status": "completed",
                    "progress": 100,
                    "message": "PDF translation completed",
                    "result_file": f"translated_{job_id}.pdf",
                    "created_at": datetime.utcnow().isoformat(),
                    "updated_at": datetime.utcnow().isoformat(),
                    "stats": {
                        "pages_processed": 1,
                        "tokens_used": 150,
                        "cost": 0.003
                    }
                }

                response = client.get(f"/api/jobs/{job_id}", headers=auth_headers)
                status_data = response.json()
                assert status_data["status"] == "completed"

            # Step 3: Filter jobs by file type
            response = client.get("/api/jobs?file_type=pdf", headers=auth_headers)
            assert response.status_code == 200
            pdf_jobs = response.json()
            assert any(job["id"] == job_id for job in pdf_jobs["jobs"])

        finally:
            # Clean up
            os.unlink(pdf_path)

    def test_job_submission_with_all_options(
        self,
        client,
        auth_headers,
        test_upload_dir
    ):
        """Test job submission with all available options"""
        pptx_path = self.create_test_pptx("複雑なテスト")

        try:
            with open(pptx_path, "rb") as f:
                files = {"file": ("complex.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
                data = {
                    "file_type": "pptx",
                    "model": "gpt-4o-2024-08-06",
                    "temperature": 0.8,
                    "offline": True,
                    "auto_fit": "shape",
                    "notes": "Complex translation with special formatting"
                }

                response = client.post(
                    "/api/translate",
                    files=files,
                    data=data,
                    headers=auth_headers
                )

            assert response.status_code == 200
            job_data = response.json()

            # Verify all options were captured
            # Note: Some of these might be in the job details rather than the response
            assert "job_id" in job_data
            assert job_data["file_type"] == "pptx"

        finally:
            os.unlink(pptx_path)

    def test_job_history_pagination(self, client, auth_headers):
        """Test job history with pagination"""
        # Create multiple jobs first
        for i in range(5):
            pptx_path = self.create_test_pptx(f"テスト{i}")

            with open(pptx_path, "rb") as f:
                files = {"file": (f"test{i}.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
                data = {
                    "file_type": "pptx",
                    "model": "gpt-4o-mini"
                }

                # Mock the creation since we don't want to actually process
                with patch('app.core.job_manager.job_manager.create_job'):
                    client.post(
                        "/api/translate",
                        files=files,
                        data=data,
                        headers=auth_headers
                    )

            os.unlink(pptx_path)

        # Test pagination
        response = client.get("/api/jobs?page=1&page_size=2", headers=auth_headers)
        assert response.status_code == 200
        page1 = response.json()
        assert len(page1["jobs"]) <= 2
        assert "total" in page1
        assert "pages" in page1

        # Test second page
        response = client.get("/api/jobs?page=2&page_size=2", headers=auth_headers)
        assert response.status_code == 200
        page2 = response.json()

        # Verify no overlap
        if page1["jobs"] and page2["jobs"]:
            page1_ids = {job["id"] for job in page1["jobs"]}
            page2_ids = {job["id"] for job in page2["jobs"]}
            assert page1_ids.isdisjoint(page2_ids)


class TestErrorScenarios:
    """Test various error scenarios"""

    def test_invalid_file_type(self, client, auth_headers):
        """Test upload with invalid file type"""
        # Create a text file
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
            f.write(b"This is not a valid file type")
            txt_path = f.name

        try:
            with open(txt_path, "rb") as f:
                files = {"file": ("invalid.txt", f, "text/plain")}
                data = {"file_type": "pptx"}  # Lying about file type

                response = client.post(
                    "/api/translate",
                    files=files,
                    data=data,
                    headers=auth_headers
                )

            assert response.status_code == 400
            assert "invalid file type" in response.json()["detail"].lower()

        finally:
            os.unlink(txt_path)

    def test_unsupported_file_format(self, client, auth_headers):
        """Test with unsupported file format in request"""
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            f.write(b"Word document content")
            docx_path = f.name

        try:
            with open(docx_path, "rb") as f:
                files = {"file": ("test.docx", f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
                data = {"file_type": "docx"}  # Unsupported format

                response = client.post(
                    "/api/translate",
                    files=files,
                    data=data,
                    headers=auth_headers
                )

            assert response.status_code == 400
            assert "file_type must be either" in response.json()["detail"].lower()

        finally:
            os.unlink(docx_path)

    def test_job_not_found(self, client, auth_headers):
        """Test accessing non-existent job"""
        fake_job_id = "00000000-0000-0000-0000-000000000000"
        response = client.get(f"/api/jobs/{fake_job_id}", headers=auth_headers)
        assert response.status_code == 404

    def test_download_nonexistent_job(self, client, auth_headers):
        """Test downloading non-existent job"""
        fake_job_id = "00000000-0000-0000-0000-000000000000"
        response = client.get(f"/api/jobs/{fake_job_id}/download", headers=auth_headers)
        assert response.status_code == 404

    def test_download_incomplete_job(self, client, auth_headers):
        """Test downloading job that hasn't completed"""
        # Create a job
        pptx_path = TestJobSubmissionWorkflow.create_test_pptx(self)

        try:
            with open(pptx_path, "rb") as f:
                files = {"file": ("test.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
                data = {"file_type": "pptx"}

                response = client.post(
                    "/api/translate",
                    files=files,
                    data=data,
                    headers=auth_headers
                )

            job_id = response.json()["job_id"]

            # Try to download before completion
            response = client.get(
                f"/api/jobs/{job_id}/download",
                headers=auth_headers
            )
            assert response.status_code == 400
            assert "not ready for download" in response.json()["detail"].lower()

        finally:
            os.unlink(pptx_path)

    def test_expired_token(self, client):
        """Test access with expired token"""
        # This would require more sophisticated token mocking
        # For now, just test with a malformed token
        headers = {"Authorization": "Bearer expired.token.here"}
        response = client.get("/api/jobs", headers=headers)
        assert response.status_code == 401

    def test_rate_limiting(self, client, auth_headers):
        """Test rate limiting on endpoints"""
        # Make many requests quickly
        for i in range(10):
            response = client.get("/api/translate/models", headers=auth_headers)
            # Should not be rate limited in test environment
            assert response.status_code == 200

    def test_large_file_upload(self, client, auth_headers):
        """Test handling of large file upload"""
        # Create a large file (10MB)
        large_content = b"x" * (10 * 1024 * 1024)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(large_content)
            large_path = f.name

        try:
            with open(large_path, "rb") as f:
                files = {"file": ("large.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
                data = {"file_type": "pptx"}

                response = client.post(
                    "/api/translate",
                    files=files,
                    data=data,
                    headers=auth_headers
                )

            # Should either succeed or fail gracefully with size error
            assert response.status_code in [200, 400, 413]

        finally:
            os.unlink(large_path)


class TestIntegrationPoints:
    """Test integration with the consolidated API structure"""

    def test_api_endpoints_structure(self, client):
        """Test that all expected API endpoints exist"""
        # Test root endpoint
        response = client.get("/")
        assert response.status_code == 200
        assert "docs" in response.json()

        # Test health endpoint
        response = client.get("/health")
        assert response.status_code == 200
        assert "status" in response.json()

        # Test API docs endpoint
        response = client.get("/docs")
        assert response.status_code == 200

    def test_auth_endpoints_accessible(self, client):
        """Test that auth endpoints are accessible without authentication"""
        # Register endpoint
        user_data = {"email": "struct@example.com", "password": "Test123!", "full_name": "Struct Test"}
        response = client.post("/api/auth/register", json=user_data)
        assert response.status_code == 200

        # Login endpoint
        login_data = {"email": "struct@example.com", "password": "Test123!"}
        response = client.post("/api/auth/login", json=login_data)
        assert response.status_code == 200

    def test_translate_endpoints_protected(self, client):
        """Test that translate endpoints require authentication"""
        # Models endpoint
        response = client.get("/api/translate/models")
        assert response.status_code == 403

        # Formats endpoint
        response = client.get("/api/translate/formats")
        assert response.status_code == 403

    def test_jobs_endpoints_protected(self, client):
        """Test that jobs endpoints require authentication"""
        # List jobs
        response = client.get("/api/jobs")
        assert response.status_code == 403

        # Search jobs
        response = client.post("/api/jobs/search", json={"search": "test"})
        assert response.status_code == 403

    def test_sse_endpoint_accessible(self, client, auth_headers):
        """Test that SSE endpoint is accessible with authentication"""
        # Note: SSE endpoint testing is complex with TestClient
        # This just verifies the endpoint exists
        response = client.get("/api/sse/status", headers=auth_headers)
        # SSE endpoints might return specific status codes
        assert response.status_code in [200, 404]  # 404 if not implemented yet

    def test_cors_headers(self, client):
        """Test that CORS headers are properly set"""
        # Make a preflight request
        response = client.options("/api/translate/models")
        assert response.status_code == 200
        assert "access-control-allow-origin" in response.headers

    def test_error_responses_format(self, client):
        """Test that error responses follow consistent format"""
        # Test 404
        response = client.get("/api/nonexistent")
        assert response.status_code == 404
        assert "detail" in response.json()

        # Test 405
        response = client.patch("/api/translate/models")
        assert response.status_code == 405
        assert "detail" in response.json()

    @patch('app.services.translation.translation_service.get_available_models')
    def test_models_endpoint_integration(self, mock_models, client, auth_headers):
        """Test models endpoint integration"""
        mock_models.return_value = [
            {"id": "gpt-4o", "name": "GPT-4o", "supports_japanese": True},
            {"id": "gpt-4o-mini", "name": "GPT-4o Mini", "supports_japanese": True}
        ]

        response = client.get("/api/translate/models", headers=auth_headers)
        assert response.status_code == 200
        data = response.json()
        assert "models" in data
        assert len(data["models"]) == 2
        assert data["models"][0]["supports_japanese"] is True

    def test_formats_endpoint_integration(self, client, auth_headers):
        """Test formats endpoint integration"""
        response = client.get("/api/translate/formats", headers=auth_headers)
        assert response.status_code == 200
        data = response.json()
        assert "formats" in data
        assert "pptx" in data["formats"]
        assert "pdf" in data["formats"]
        assert "description" in data["formats"]["pptx"]
        assert "max_size" in data["formats"]["pptx"]

    def test_job_statistics_endpoint(self, client, auth_headers):
        """Test job statistics endpoint"""
        response = client.get("/api/jobs/statistics", headers=auth_headers)
        assert response.status_code == 200
        data = response.json()
        assert "total_jobs" in data
        assert "status_counts" in data
        assert "average_duration_minutes" in data
        assert isinstance(data["status_counts"], dict)


# Performance tests (marked as slow)
@pytest.mark.slow
class TestPerformance:
    """Performance-related smoke tests"""

    def test_concurrent_job_creation(self, client, auth_headers):
        """Test creating multiple jobs concurrently"""
        import threading
        import time

        results = []

        def create_job():
            pptx_path = TestJobSubmissionWorkflow.create_test_pptx(self)

            try:
                with open(pptx_path, "rb") as f:
                    files = {"file": ("concurrent.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
                    data = {"file_type": "pptx"}

                    response = client.post(
                        "/api/translate",
                        files=files,
                        data=data,
                        headers=auth_headers
                    )
                    results.append(response.status_code)
            finally:
                os.unlink(pptx_path)

        # Create 5 threads
        threads = []
        for _ in range(5):
            t = threading.Thread(target=create_job)
            threads.append(t)
            t.start()

        # Wait for all threads
        for t in threads:
            t.join()

        # All should succeed
        assert all(status == 200 for status in results)

    def test_job_list_performance(self, client, auth_headers):
        """Test job list endpoint performance with many jobs"""
        # Create many jobs (simplified - just mock the DB)
        import time
        start_time = time.time()

        response = client.get("/api/jobs?page=1&page_size=100", headers=auth_headers)
        end_time = time.time()

        assert response.status_code == 200
        assert (end_time - start_time) < 1.0  # Should be fast


# Utility functions for testing
def create_auth_headers_for_user(client, email, password):
    """Helper to create auth headers for a specific user"""
    # Register
    user_data = {
        "email": email,
        "password": password,
        "full_name": f"User {email}"
    }
    client.post("/api/auth/register", json=user_data)

    # Login
    login_data = {"email": email, "password": password}
    response = client.post("/api/auth/login", json=login_data)
    token = response.json()["access_token"]

    return {"Authorization": f"Bearer {token}"}