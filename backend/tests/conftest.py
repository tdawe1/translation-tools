import os
import tempfile
import shutil
import sys
import pytest
from fastapi.testclient import TestClient
import warnings
from pathlib import Path
from datetime import datetime

# Set up test environment BEFORE any imports
# Set pytest marker for test environment detection
os.environ["PYTEST_RUNNING"] = "1"

# Load test environment variables from .env.test if it exists
env_test_path = os.path.join(os.path.dirname(__file__), '..', '.env.test')
if os.path.exists(env_test_path):
    with open(env_test_path, 'r') as f:
        for line in f:
            line = line.strip()
            if line and not line.startswith('#') and '=' in line:
                key, value = line.split('=', 1)
                value = value.split('#')[0].strip()
                key = key.strip()
                os.environ[key] = value

# Override specific environment variables for testing (these take precedence over .env.test)
os.environ["DEBUG"] = "true"
os.environ["SECRET_KEY"] = "test-secret-key-for-pytest-testing-only-32-chars-long"
os.environ["OPENAI_API_KEY"] = "mock-sk-for-testing"
os.environ["DATABASE_URL"] = "sqlite:///:memory:"
os.environ["LOG_LEVEL"] = "WARNING"
os.environ["UPLOAD_DIR"] = "test_uploads"
os.environ["OUTPUT_DIR"] = "test_outputs"

# Add the app directory to the Python path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

# Now we can safely import app modules
from app.main import app
from app.core.config import settings
from app.database.session import get_db
from app.database.database import Base
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker, Session


@pytest.fixture(scope="session", autouse=True)
def test_settings():
    """Configure test settings that apply to all tests."""
    # Suppress warnings about test keys
    warnings.filterwarnings("ignore", message="DEBUG mode.*SECRET_KEY.*", category=UserWarning)
    warnings.filterwarnings("ignore", message="DEBUG mode.*OPENAI_API_KEY.*", category=UserWarning)

    # Ensure test environment variables are set
    assert os.environ.get("DEBUG") == "true", "DEBUG must be true for tests"
    assert os.environ.get("DATABASE_URL").startswith("sqlite"), "Must use SQLite for tests"
    assert "test" in os.environ.get("UPLOAD_DIR", ""), "Must use test upload directory"
    assert "test" in os.environ.get("OUTPUT_DIR", ""), "Must use test output directory"


@pytest.fixture(scope="function", autouse=True)
def clean_test_environment():
    """Ensure a clean test environment before each test."""
    # Clean test directories if they exist
    for dir_name in ["test_uploads", "test_outputs"]:
        dir_path = Path(dir_name)
        if dir_path.exists():
            shutil.rmtree(dir_path, ignore_errors=True)
        # Ensure directories exist
        dir_path.mkdir(exist_ok=True)

    yield

    # Clean up test directories after test
    for dir_name in ["test_uploads", "test_outputs"]:
        dir_path = Path(dir_name)
        if dir_path.exists():
            shutil.rmtree(dir_path, ignore_errors=True)


@pytest.fixture(scope="function")
def test_upload_dir():
    """Create a temporary upload directory for testing."""
    temp_upload = tempfile.mkdtemp()

    # Override settings temporarily
    original_upload = settings.UPLOAD_DIR
    settings.UPLOAD_DIR = temp_upload

    yield temp_upload

    # Cleanup
    shutil.rmtree(temp_upload, ignore_errors=True)
    settings.UPLOAD_DIR = original_upload


@pytest.fixture(scope="function")
def test_output_dir():
    """Create a temporary output directory for testing."""
    temp_output = tempfile.mkdtemp()

    # Override settings temporarily
    original_output = settings.OUTPUT_DIR
    settings.OUTPUT_DIR = temp_output

    yield temp_output

    # Cleanup
    shutil.rmtree(temp_output, ignore_errors=True)
    settings.OUTPUT_DIR = original_output


@pytest.fixture(scope="function")
def test_db():
    """Create a test database session."""
    # Use in-memory SQLite for tests
    test_engine = create_engine("sqlite:///:memory:", connect_args={"check_same_thread": False})
    TestSessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=test_engine)

    # Create tables
    Base.metadata.create_all(bind=test_engine)

    # Override the database dependency
    def override_get_db():
        try:
            db = TestSessionLocal()
            yield db
        finally:
            db.close()

    app.dependency_overrides[get_db] = override_get_db

    # Return a session for the test
    db_session = TestSessionLocal()
    try:
        yield db_session
    finally:
        db_session.close()
        app.dependency_overrides.clear()
        Base.metadata.drop_all(bind=test_engine)
        test_engine.dispose()


@pytest.fixture(scope="function")
def client(test_db):
    """Create a test client with test settings."""
    # The test_db fixture already sets up the database override
    with TestClient(app) as test_client:
        yield test_client


@pytest.fixture(scope="function")
def auth_headers(client):
    """Create authentication headers for protected endpoints."""
    # Register a test user
    user_data = {
        "email": "test@example.com",
        "password": "testpassword123",
        "full_name": "Test User"
    }
    client.post("/api/auth/register", json=user_data)

    # Login to get token
    login_data = {
        "email": "test@example.com",
        "password": "testpassword123"
    }
    response = client.post("/api/auth/login", json=login_data)
    token = response.json()["access_token"]

    return {"Authorization": f"Bearer {token}"}


@pytest.fixture(scope="function")
def admin_auth_headers(client):
    """Create admin authentication headers."""
    # Register an admin user (this assumes you have a way to create admins)
    admin_data = {
        "email": "admin@example.com",
        "password": "adminpassword123",
        "full_name": "Admin User"
    }
    client.post("/api/auth/register", json=admin_data)

    # Login to get token
    login_data = {
        "email": "admin@example.com",
        "password": "adminpassword123"
    }
    response = client.post("/api/auth/login", json=login_data)
    token = response.json()["access_token"]

    return {"Authorization": f"Bearer {token}"}


@pytest.fixture(scope="function")
def mock_openai(monkeypatch):
    """Mock OpenAI API calls to avoid real API calls in tests."""
    class MockOpenAI:
        def __init__(self, *args, **kwargs):
            pass

        class chat:
            class completions:
                def create(self, *args, **kwargs):
                    # Return a mock translation response
                    return {
                        "choices": [{
                            "message": {
                                "content": '{"translations": [{"original": "テスト", "translated": "Test"}]}'
                            }
                        }]
                    }

    # Mock the OpenAI client
    monkeypatch.setattr("openai.OpenAI", MockOpenAI)

    # Also mock any other OpenAI-related imports
    monkeypatch.setattr("app.services.translation.openai", MockOpenAI)

    return MockOpenAI


@pytest.fixture(scope="function")
def sample_pptx_file():
    """Create a minimal PPTX file for testing."""
    try:
        from pptx import Presentation

        # Create a simple presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "テスト"
        slide.placeholders[1].text = "Test content"

        # Save to temporary file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            yield f.name

        # Cleanup
        os.unlink(f.name)

    except ImportError:
        # If python-pptx is not available, create a minimal mock file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            # Write minimal valid PPTX structure
            f.write(b"PK\x03\x04")  # ZIP header
            f.write(b"MOCK_PPTX_FILE" * 100)  # Mock content
            yield f.name
        os.unlink(f.name)


@pytest.fixture(scope="function")
def sample_pdf_file():
    """Create a minimal PDF file for testing."""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter

        # Create a simple PDF
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            c = canvas.Canvas(f.name, pagesize=letter)
            c.drawString(100, 750, "テスト")
            c.drawString(100, 730, "Test PDF content")
            c.save()
            yield f.name

        # Cleanup
        os.unlink(f.name)

    except ImportError:
        # If reportlab is not available, create a minimal mock file
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(b"%PDF-1.4\n")  # PDF header
            f.write(b"MOCK_PDF_FILE" * 100)  # Mock content
            yield f.name
        os.unlink(f.name)


@pytest.fixture(scope="function")
def mock_job_manager(monkeypatch):
    """Mock the job manager to avoid actual job processing."""
    class MockJobManager:
        def create_job(self, *args, **kwargs):
            return {
                "id": "test-job-id",
                "status": "pending",
                "created_at": datetime.utcnow().isoformat()
            }

        def get_job_status(self, job_id):
            return {
                "id": job_id,
                "status": "completed",
                "progress": 100,
                "message": "Translation completed",
                "result_file": f"translated_{job_id}.pptx"
            }

        def process_translation_job(self, job_id):
            pass  # Mock processing

    mock_manager = MockJobManager()
    monkeypatch.setattr("app.core.job_manager.job_manager", mock_manager)
    return mock_manager