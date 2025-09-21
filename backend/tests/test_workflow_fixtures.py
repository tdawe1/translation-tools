"""
Comprehensive test fixtures and utilities for workflow testing.
"""
import pytest
import os
import tempfile
import json
from pathlib import Path
from unittest.mock import patch, MagicMock
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.colors import black, blue, red
from reportlab.lib.units import inch
from datetime import datetime, timedelta
import zipfile
import xml.etree.ElementTree as ET
from io import BytesIO


@pytest.fixture(scope="session")
def sample_pptx_with_japanese():
    """Create a sample PPTX file with Japanese text for testing"""
    prs = Presentation()

    # Title slide
    title_slide = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide)
    slide1.shapes.title.text = "日本語プレゼンテーション"
    slide1.placeholders[1].text = "テスト用のサンプルファイル"

    # Content slide with multiple text boxes
    content_slide = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_slide)
    slide2.shapes.title.text = "プロジェクト概要"

    # Add text box with Japanese content
    left = inch
    top = inch * 2
    width = inch * 6
    height = inch * 4
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "これは翻訳テストのためのサンプルテキストです。"

    # Add bullet points
    p = tf.add_paragraph()
    p.text = "プロジェクトの目標"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "技術的な要件"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "スケジュールと予算"
    p.level = 1

    # Table slide
    table_slide = prs.slide_layouts[5]
    slide3 = prs.slides.add_slide(table_slide)
    slide3.shapes.title.text = "プロジェクトデータ"

    # Add a table
    rows, cols = 3, 2
    left = inch
    top = inch * 2
    width = inch * 6
    height = inch * 2
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

    # Fill table
    table.cell(0, 0).text = "項目"
    table.cell(0, 1).text = "値"
    table.cell(1, 0).text = "予算"
    table.cell(1, 1).text = "￥10,000,000"
    table.cell(2, 0).text = "期間"
    table.cell(2, 1).text = "6ヶ月"

    # Save to temporary file
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        yield f.name

    # Cleanup
    os.unlink(f.name)


@pytest.fixture(scope="session")
def sample_pdf_with_japanese():
    """Create a sample PDF file with Japanese text for testing"""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # Try to use a Japanese font if available
    japanese_fonts = [
        "/usr/share/fonts/truetype/vlgothic/VL-Gothic-Regular.ttf",
        "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
        "/usr/share/fonts/truetype/takao-gothic/TakaoGothic.ttf"
    ]

    font_registered = False
    for font_path in japanese_fonts:
        if os.path.exists(font_path):
            try:
                pdfmetrics.registerFont(TTFont('JapaneseFont', font_path))
                font_registered = True
                break
            except:
                continue

    # Create PDF
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
        c = canvas.Canvas(f.name, pagesize=A4)

        # Set font
        if font_registered:
            c.setFont('JapaneseFont', 12)
        else:
            c.setFont('Helvetica', 12)

        # Title
        c.drawString(72, 750, "Japanese Document Test")
        c.drawString(72, 730, "日本語文書テスト")

        # Content
        y = 700
        lines = [
            "これはPDF翻訳テストのためのサンプルテキストです。",
            "複数行にわたるコンテンツを含んでいます。",
            "数字、日付、特殊文字も含めてテストします。",
            "Price: ￥5,000 (税込)",
            "Date: 2024年1月15日",
            "Contact: test@example.com"
        ]

        for line in lines:
            c.drawString(72, y, line)
            y -= 20

        # Add a simple table using lines
        y -= 20
        c.drawString(72, y, "Table Test:")
        y -= 20
        c.line(72, y, 300, y)  # Header line
        y -= 15
        c.drawString(72, y, "Item")
        c.drawString(200, y, "Value")
        y -= 15
        c.line(72, y, 300, y)  # Separator
        y -= 15
        c.drawString(72, y, "Product A")
        c.drawString(200, y, "￥2,500")
        y -= 15
        c.line(72, y, 300, y)  # Footer

        c.showPage()
        c.save()
        yield f.name

    # Cleanup
    os.unlink(f.name)


@pytest.fixture
def mock_translation_responses():
    """Mock translation API responses for testing"""
    return {
        "success": {
            "choices": [{
                "message": {
                    "content": json.dumps({
                        "translations": [
                            {"original": "テスト", "translated": "Test"},
                            {"original": "日本語", "translated": "Japanese"},
                            {"original": "翻訳", "translated": "Translation"}
                        ]
                    })
                }
            }]
        },
        "error": {
            "error": {
                "message": "API error occurred",
                "type": "invalid_request_error"
            }
        },
        "partial": {
            "choices": [{
                "message": {
                    "content": json.dumps({
                        "translations": [
                            {"original": "テスト", "translated": "Test"},
                            {"original": "日本語", "translated": ""}  # Empty translation
                        ]
                    })
                }
            }]
        }
    }


@pytest.fixture
def mock_job_database():
    """Mock job database entries for testing"""
    return [
        {
            "id": "550e8400-e29b-41d4-a716-446655440000",
            "user_id": "user123",
            "file_name": "test1.pptx",
            "file_type": "pptx",
            "status": "completed",
            "created_at": datetime.utcnow().isoformat(),
            "updated_at": datetime.utcnow().isoformat(),
            "progress": 100,
            "message": "Translation completed",
            "input_file": "uploads/test1.pptx",
            "output_file": "outputs/test1_en.pptx",
            "config": {
                "model": "gpt-4o-mini",
                "temperature": 0.6,
                "auto_fit": "norm"
            },
            "stats": {
                "pages_processed": 5,
                "tokens_used": 1500,
                "cost": 0.03
            }
        },
        {
            "id": "550e8400-e29b-41d4-a716-446655440001",
            "user_id": "user123",
            "file_name": "test2.pdf",
            "file_type": "pdf",
            "status": "processing",
            "created_at": datetime.utcnow().isoformat(),
            "updated_at": datetime.utcnow().isoformat(),
            "progress": 45,
            "message": "Processing page 3 of 7",
            "input_file": "uploads/test2.pdf",
            "output_file": None,
            "config": {
                "model": "gpt-4o",
                "temperature": 0.7,
                "pages": "1-10"
            },
            "stats": {
                "pages_processed": 3,
                "tokens_used": 800,
                "cost": 0.016
            }
        }
    ]


@pytest.fixture
def sample_user_profiles():
    """Sample user profiles for testing"""
    return [
        {
            "id": "user123",
            "email": "test@example.com",
            "full_name": "Test User",
            "is_active": True,
            "created_at": datetime.utcnow().isoformat(),
            "preferences": {
                "default_model": "gpt-4o-mini",
                "default_temperature": 0.6,
                "notifications_enabled": True
            }
        },
        {
            "id": "user456",
            "email": "admin@example.com",
            "full_name": "Admin User",
            "is_active": True,
            "is_admin": True,
            "created_at": datetime.utcnow().isoformat(),
            "preferences": {
                "default_model": "gpt-4o",
                "default_temperature": 0.7,
                "notifications_enabled": True
            }
        }
    ]


@pytest.fixture
def test_auth_scenarios():
    """Test scenarios for authentication edge cases"""
    return {
        "expired_token": "eyJ0eXAiOiJKV1QiLCJhbGciOiJIUzI1NiJ9.eyJzdWIiOiJ0ZXN0QGV4YW1wbGUuY29tIiwiZXhwIjoxNjAwMDAwMDAwfQ.invalid_signature",
        "malformed_token": "this.is.not.a.valid.jwt",
        "empty_token": "",
        "none_token": None
    }


# Helper functions for test data creation
def create_test_pptx_with_slides(num_slides=3):
    """Create a test PPTX with specified number of slides"""
    prs = Presentation()

    for i in range(num_slides):
        slide_layout = prs.slide_layouts[1] if i > 0 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        if i == 0:
            slide.shapes.title.text = f"Test Presentation ({num_slides} slides)"
            slide.placeholders[1].text = "Generated for testing purposes"
        else:
            slide.shapes.title.text = f"Slide {i}"
            slide.placeholders[1].text = f"This is content for slide {i}.\n" * 5

    # Save to BytesIO for in-memory use
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def create_test_pdf_with_pages(num_pages=3):
    """Create a test PDF with specified number of pages"""
    buffer = BytesIO()
    p = canvas.Canvas(buffer, pagesize=A4)

    for i in range(num_pages):
        p.drawString(72, 750, f"Test PDF - Page {i+1}")
        p.drawString(72, 720, f"This is page {i+1} of {num_pages}")
        p.drawString(72, 690, "Content for testing translation services.")

        # Add some Japanese text
        p.drawString(72, 660, f"これはテストページ{i+1}です。")

        p.showPage()

    p.save()
    buffer.seek(0)
    return buffer


def create_mock_translation_result(texts):
    """Create a mock translation result"""
    translations = []
    for text in texts:
        translations.append({
            "original": text,
            "translated": f"[Translated] {text}"
        })
    return {"translations": translations}


def create_test_user_data(overrides=None):
    """Create test user data with optional overrides"""
    data = {
        "email": "test@example.com",
        "password": "TestPassword123!",
        "full_name": "Test User"
    }
    if overrides:
        data.update(overrides)
    return data


def create_test_job_data(overrides=None):
    """Create test job data with optional overrides"""
    data = {
        "file_type": "pptx",
        "model": "gpt-4o-mini",
        "temperature": 0.6,
        "auto_fit": "norm",
        "offline": False
    }
    if overrides:
        data.update(overrides)
    return data


# Test data generators
@pytest.fixture
def test_data_generator():
    """Fixture providing test data generation methods"""
    class TestDataGenerator:
        @staticmethod
        def generate_japanese_texts(count=10):
            """Generate sample Japanese texts for testing"""
            texts = [
                "これはサンプルテキストです。",
                "日本語の翻訳をテストしています。",
                "プロジェクトの進捗状況",
                "会議の議事録",
                "技術的な仕様書",
                "ユーザーマニュアル",
                "プレゼンテーションスライド",
                "ビジネスレポート",
                "製品の説明書",
                "契約書の草案"
            ]
            return texts[:count]

        @staticmethod
        def generate_english_texts(count=10):
            """Generate sample English texts for testing"""
            texts = [
                "This is a sample text.",
                "Testing Japanese translation.",
                "Project progress status",
                "Meeting minutes",
                "Technical specifications",
                "User manual",
                "Presentation slides",
                "Business report",
                "Product description",
                "Contract draft"
            ]
            return texts[:count]

        @staticmethod
        def generate_mixed_content():
            """Generate content with mixed Japanese and English"""
            return [
                "プロジェクト名: Project Alpha",
                "期限: Deadline 2024-12-31",
                "担当者: Responsible - Yamada Taro",
                "ステータス: Status - In Progress",
                "予算: Budget ￥5,000,000"
            ]

        @staticmethod
        def generate_special_cases():
            """Generate special translation cases"""
            return [
                "123-456-7890",  # Numbers
                "test@example.com",  # Email
                "https://example.com",  # URL
                "￥10,000",  # Currency
                "2024年1月15日",  # Date
                "株式会社テスト",  # Company
                "山田太郎様",  # Name with honorific
                "〒100-0001 東京都",  # Address
                "TEL: 03-1234-5678",  # Phone
                "注意: 注意書き"  # Warning
            ]

    return TestDataGenerator()