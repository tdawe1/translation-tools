#!/usr/bin/env python3
"""
English expansion policy: compress → spill → shrink
Handles text overflow in translated slides systematically.
"""
import re
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE

def calculate_expansion_ratio(original_jp: str, translated_en: str) -> float:
    """Calculate expansion ratio between Japanese and English text."""
    jp_len = len(original_jp.strip())
    en_len = len(translated_en.strip())
    return en_len / jp_len if jp_len > 0 else 1.0

def condense_block(client, tagged_en: str, target_ratio: float = 0.85) -> str:
    """Stage 1: Compress text by removing filler while preserving meaning."""
    reduction_pct = int((1 - target_ratio) * 100)
    prompt = f"""Shorten this English slide text by ~{reduction_pct}%.

KEY REQUIREMENTS:
- Preserve exact tags [b][i][u][sup][sub][li-lN]…[/li] and placeholders ⟦…⟧
- Keep all numbers, URLs, and specific terms intact
- Keep bullets as concise fragments, not full sentences
- Remove filler words: "in order to"→"to", "utilize"→"use", "as well as"→"and"
- Drop articles where clear ("the", "a") and most instances of "that"
- Use one verb per bullet; cut unnecessary adverbs
- Collapse double nouns: "customer onboarding process"→"customer onboarding"
- Maintain parallel structure and professional tone
- Do not add or remove actual content/meaning

Text to shorten:"""

    try:
        resp = client.responses.create(
            model="gpt-5",
            reasoning_effort="high",  # Think carefully about what to cut
            verbosity="low",
            input=[{"role": "user", "content": prompt + "\n\n" + tagged_en}],
            response_format={"type": "text"},
            temperature=0.2,  # Conservative for preservation
        )
        return resp.output_text().strip()
    except Exception:
        return tagged_en  # Fallback to original if compression fails

def spill_to_notes(text_block: str, overflow_content: str) -> tuple[str, str]:
    """Stage 2: Move overflow content to slide notes with reference stub."""
    # Find the last complete sentence or bullet point that fits
    sentences = re.split(r'(?<=[.!?])\s+', text_block)
    if len(sentences) <= 1:
        return text_block, overflow_content
    
    # Keep most content, add reference to notes
    main_content = sentences[:-1]
    spilled_content = sentences[-1]
    
    stub_text = " ".join(main_content) + " (details → Notes)"
    notes_content = f"Additional details:\n{spilled_content}\n{overflow_content}"
    
    return stub_text, notes_content

def tighten_textframe(tf, is_title: bool = False):
    """Stage 3: Apply layout optimizations to buy space."""
    # Enable shrink-to-fit autofit
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    # Tighten margins
    try:
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
    except Exception:
        pass

    tf.word_wrap = True

    # Paragraph-level spacing and bullet indents
    for p in tf.paragraphs:
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        try:
            p.line_spacing = 1.1  # Tighter line spacing
        except Exception:
            pass

        # Tighten bullet indents
        if getattr(p, "level", 0) > 0:
            if p.level == 1:
                p.left_indent = Inches(0.3)
                p.hanging_indent = Inches(0.15)
            else:
                p.left_indent = Inches(0.25)
                p.hanging_indent = Inches(0.12)
        else:
            p.left_indent = Inches(0.0)
            p.hanging_indent = Inches(0.18)

        # Enforce minimum font sizes
        for r in p.runs:
            if r.font.size:
                min_pt = Pt(18) if is_title else Pt(11)
                if r.font.size < min_pt:
                    r.font.size = min_pt

def apply_expansion_policy(client, slide_content: dict, original_jp: str) -> dict:
    """Apply 3-stage expansion policy based on content type and expansion ratio."""
    translated_en = slide_content.get("translated_text", "")
    content_type = slide_content.get("content_type", "body")  # title, body, table
    
    expansion_ratio = calculate_expansion_ratio(original_jp, translated_en)
    
    # Define thresholds by content type
    thresholds = {
        "title": 1.8,
        "body": 1.4, 
        "table": 1.2,
    }
    
    threshold = thresholds.get(content_type, 1.4)
    
    if expansion_ratio <= threshold:
        return slide_content  # No intervention needed
    
    # Stage 1: Compress
    condensed_text = condense_block(client, translated_en, target_ratio=0.85)
    new_ratio = calculate_expansion_ratio(original_jp, condensed_text)
    
    result = slide_content.copy()
    result["translated_text"] = condensed_text
    result["expansion_ratio"] = new_ratio
    result["compression_applied"] = True
    
    # Stage 2: Spill to notes if still too long
    if new_ratio > threshold:
        stub_text, notes_content = spill_to_notes(condensed_text, "")
        result["translated_text"] = stub_text
        result["notes_content"] = notes_content
        result["spill_applied"] = True
        
    # Stage 3: Layout tightening will be applied during PPTX write-back
    result["needs_tightening"] = new_ratio > (threshold * 0.9)
    
    return result

def get_style_guide_additions() -> str:
    """Return style guide additions for consistent shortening."""
    return """
CONCISENESS RULES for slide translation:
- Use fragments, not full sentences in bullets
- Remove filler: "in order to"→"to", "utilize"→"use", "as well as"→"and"  
- Drop articles where clear: "the", "a"
- Cut most instances of "that"
- Use symbols: "and"→"&" in labels, "approximately"→"~", "versus"→"vs."
- One verb per bullet; cut adverbs
- Collapse double nouns: "customer onboarding process"→"customer onboarding"
- Keep parallel structure in bullet lists
"""

# Integration points for main translator
def should_use_expansion_policy(model: str) -> bool:
    """Check if we should apply expansion policy based on model capabilities."""
    return model.startswith("gpt-5")  # Only for models that support reasoning

def integrate_with_batch_translate():
    """
    Integration example for translate_pptx_inplace.py:
    
    # After translation but before write-back:
    for slide_data in translated_slides:
        if should_use_expansion_policy(model):
            slide_data = apply_expansion_policy(client, slide_data, original_jp_text)
            
        # During PPTX write-back:
        if slide_data.get("needs_tightening", False):
            for text_frame in slide.text_frames:
                is_title = (text_frame.parent_shape.name or "").lower().startswith("title")
                tighten_textframe(text_frame, is_title)
    """
    pass